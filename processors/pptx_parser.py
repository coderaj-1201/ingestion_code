"""
PPTX Parser
===========

PowerPoint slides map naturally to chunks:
  - Slide title → section_heading, chunk_type=heading
  - Slide body text → paragraph chunks
  - Tables on slide → table chunks with NL summary
  - Slide number → page_number
  - Presentation title from core properties → doc_title

Notes: slide notes are ignored (usually presenter notes, not policy content).
"""
from __future__ import annotations

import logging
from datetime import datetime, timezone
from uuid import uuid4

from pptx import Presentation
from pptx.util import Pt

from processors.pdf_parser import _llm_clean_page, _llm_serialise_table
from shared.models import ChunkType, RawChunk

logger = logging.getLogger(__name__)


def _shape_is_title(shape) -> bool:
    from pptx.enum.shapes import PP_PLACEHOLDER
    try:
        return shape.placeholder_format and shape.placeholder_format.idx in (0, 1)
    except Exception:
        return False


def _table_to_markdown_pptx(table) -> str:
    rows = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(rows)


def parse_pptx(
    file_bytes: bytes,
    doc_name: str,
    doc_url: str,
    domain: str,
    blob_path: str,
) -> list[RawChunk]:
    import io
    ingested_at = datetime.now(timezone.utc).isoformat()
    prs         = Presentation(io.BytesIO(file_bytes))

    doc_title = doc_name.replace(".pptx", "").replace("_", " ")
    if prs.core_properties and prs.core_properties.title:
        doc_title = prs.core_properties.title

    chunks: list[RawChunk] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_title  = ""
        slide_body:  list[str] = []
        parent_id    = str(uuid4())
        # Collect all child chunks for this slide before appending to `chunks`.
        # This ensures the parent chunk is always written first (ordering guarantee
        # required by the embedding agent when it uploads parents before children).
        slide_chunks: list[RawChunk] = []

        # Extract title and body text from shapes
        for shape in slide.shapes:
            if not shape.has_text_frame and not shape.has_table:
                continue

            if _shape_is_title(shape) and shape.has_text_frame:
                slide_title = shape.text_frame.text.strip()
                continue

            if shape.has_table:
                tbl_md = _table_to_markdown_pptx(shape.table)
                if tbl_md:
                    nl = _llm_serialise_table(tbl_md, slide_title)
                    slide_chunks.append(RawChunk(
                        chunk_id           = str(uuid4()),
                        parent_id          = parent_id,
                        chunk_type         = ChunkType.TABLE,
                        domain             = domain,
                        doc_name           = doc_name,
                        source             = doc_name,
                        doc_url            = doc_url,
                        file_type          = "pptx",
                        blob_path          = blob_path,
                        ingested_at        = ingested_at,
                        page_number        = slide_num,
                        title              = doc_title,
                        section_heading    = slide_title,
                        section_subheading = "",
                        content            = nl,
                        table_raw          = tbl_md,
                    ))
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_body.append(text)

        all_text = (slide_title + "\n" + "\n".join(slide_body)).strip()
        if not all_text:
            # Skip empty slides entirely — do not commit any collected slide_chunks
            # since their parent would never be written.
            continue

        # Parent first, then all children — ordering is required by the embedding agent.
        chunks.append(RawChunk(
            chunk_id           = parent_id,
            parent_id          = "",
            chunk_type         = ChunkType.HEADING if slide_title else ChunkType.PARAGRAPH,
            domain             = domain,
            doc_name           = doc_name,
            source             = doc_name,
            doc_url            = doc_url,
            file_type          = "pptx",
            blob_path          = blob_path,
            ingested_at        = ingested_at,
            page_number        = slide_num,
            title              = doc_title,
            section_heading    = slide_title,
            section_subheading = "",
            content            = all_text,
        ))

        # Table chunks collected above (they now have a committed parent)
        chunks.extend(slide_chunks)

        if slide_title:
            chunks.append(RawChunk(
                chunk_id           = str(uuid4()),
                parent_id          = parent_id,
                chunk_type         = ChunkType.HEADING,
                domain             = domain,
                doc_name           = doc_name,
                source             = doc_name,
                doc_url            = doc_url,
                file_type          = "pptx",
                blob_path          = blob_path,
                ingested_at        = ingested_at,
                page_number        = slide_num,
                title              = doc_title,
                section_heading    = slide_title,
                section_subheading = "",
                content            = slide_title,
            ))

        for body_text in slide_body:
            cleaned = _llm_clean_page(body_text, slide_num)
            if not cleaned:
                continue
            chunks.append(RawChunk(
                chunk_id           = str(uuid4()),
                parent_id          = parent_id,
                chunk_type         = ChunkType.PARAGRAPH,
                domain             = domain,
                doc_name           = doc_name,
                source             = doc_name,
                doc_url            = doc_url,
                file_type          = "pptx",
                blob_path          = blob_path,
                ingested_at        = ingested_at,
                page_number        = slide_num,
                title              = doc_title,
                section_heading    = slide_title,
                section_subheading = "",
                content            = cleaned,
            ))

    logger.info("PPTX parsed: %s → %d chunks", doc_name, len(chunks))
    return chunks
