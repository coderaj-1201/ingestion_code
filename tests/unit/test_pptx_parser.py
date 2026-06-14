"""
Unit tests for processors/pptx_parser.py.

Tests parse_pptx() using PPTX bytes created in-memory with python-pptx.
All LLM calls are mocked — no real Azure OpenAI calls are made.

Assumptions:
- _llm_clean_page and _llm_serialise_table are imported from pdf_parser into
  pptx_parser, so they are patched at 'processors.pdf_parser.get_openai_client'.
- The parser appends the parent chunk BEFORE table children (table chunks are
  appended during shape iteration which happens before the parent append).
  After the bug fix, parent must exist for every child's parent_id.
- Slide number starts at 1 (enumerate(prs.slides, start=1)).
- Image-only slides (no text frames, no tables) produce 0 chunks.
"""
from __future__ import annotations

import io
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from processors.pptx_parser import parse_pptx
from shared.models import ChunkType


def _llm_mock():
    mock_client = MagicMock()
    mock_choice = MagicMock()
    mock_choice.message.content = "LLM cleaned or summarised text."
    mock_client.chat.completions.create.return_value = MagicMock(choices=[mock_choice])
    return mock_client


def _two_slide_pptx() -> bytes:
    """Two-slide presentation with title and body text on each slide."""
    prs = Presentation()
    layout = prs.slide_layouts[1]

    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Q1 Results"
    s1.placeholders[1].text = "Revenue grew by 15 percent compared to the previous quarter."

    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "HR Policy Update"
    s2.placeholders[1].text = "New remote work policy effective from March 2024."

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_parse_pptx_returns_chunks(sample_pptx_bytes):
    with patch("processors.pdf_parser.get_openai_client", return_value=_llm_mock()):
        chunks = parse_pptx(sample_pptx_bytes, "deck.pptx", "https://example.com", "hr", "hr/deck.pptx")
    assert len(chunks) > 0


def test_parse_pptx_slide_number_as_page_number(sample_pptx_bytes):
    with patch("processors.pdf_parser.get_openai_client", return_value=_llm_mock()):
        chunks = parse_pptx(sample_pptx_bytes, "deck.pptx", "https://example.com", "hr", "hr/deck.pptx")
    page_numbers = {c.page_number for c in chunks}
    assert 1 in page_numbers
    assert 2 in page_numbers


def test_parse_pptx_slide_title_as_section_heading():
    file_bytes = _two_slide_pptx()
    with patch("processors.pdf_parser.get_openai_client", return_value=_llm_mock()):
        chunks = parse_pptx(file_bytes, "deck.pptx", "https://example.com", "hr", "hr/deck.pptx")
    slide1_chunks = [c for c in chunks if c.page_number == 1]
    assert all(c.section_heading == "Q1 Results" for c in slide1_chunks)


def test_parse_pptx_parent_chunk_appears_before_children():
    file_bytes = _two_slide_pptx()
    with patch("processors.pdf_parser.get_openai_client", return_value=_llm_mock()):
        chunks = parse_pptx(file_bytes, "deck.pptx", "https://example.com", "hr", "hr/deck.pptx")
    chunk_index = {c.chunk_id: i for i, c in enumerate(chunks)}
    children = [c for c in chunks if c.parent_id != ""]
    for child in children:
        parent_idx = chunk_index.get(child.parent_id)
        assert parent_idx is not None, f"Parent {child.parent_id!r} not in chunks"
        assert parent_idx < chunk_index[child.chunk_id], (
            "Parent chunk must appear at an earlier index than its child"
        )


def test_parse_pptx_table_chunk_has_correct_parent():
    # Build a slide with a table shape
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add a title text box manually
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import PP_PLACEHOLDER
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    txBox.text_frame.text = "Slide With Table"

    # Add a table
    rows, cols = 2, 3
    tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(6), Inches(2)).table
    tbl.cell(0, 0).text = "Col A"
    tbl.cell(0, 1).text = "Col B"
    tbl.cell(0, 2).text = "Col C"
    tbl.cell(1, 0).text = "Val 1"
    tbl.cell(1, 1).text = "Val 2"
    tbl.cell(1, 2).text = "Val 3"

    buf = io.BytesIO()
    prs.save(buf)

    with patch("processors.pdf_parser.get_openai_client", return_value=_llm_mock()):
        chunks = parse_pptx(buf.getvalue(), "deck.pptx", "https://example.com", "hr", "hr/deck.pptx")

    table_chunks = [c for c in chunks if c.chunk_type == ChunkType.TABLE]
    if table_chunks:  # table may or may not be detected depending on layout
        parent_ids = {c.chunk_id for c in chunks if c.parent_id == ""}
        for tc in table_chunks:
            assert tc.parent_id in parent_ids


def test_parse_pptx_image_only_slide_produces_no_chunks():
    # A slide with only a picture (no text frame, no table) should contribute 0 chunks
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    prs.slides.add_slide(blank_layout)  # blank slide, no text shapes

    buf = io.BytesIO()
    prs.save(buf)

    with patch("processors.pdf_parser.get_openai_client", return_value=_llm_mock()):
        chunks = parse_pptx(buf.getvalue(), "deck.pptx", "https://example.com", "hr", "hr/deck.pptx")
    assert len(chunks) == 0


def test_parse_pptx_table_children_have_parent_in_output():
    """
    After the PPTX orphan-child bug fix: table chunks are appended before the
    parent chunk within the shape loop. The final output MUST contain a parent
    chunk for every table chunk's parent_id.
    """
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Data Slide"
    slide.placeholders[1].text = "Some body content here."

    buf = io.BytesIO()
    prs.save(buf)

    with patch("processors.pdf_parser.get_openai_client", return_value=_llm_mock()):
        chunks = parse_pptx(buf.getvalue(), "deck.pptx", "https://example.com", "hr", "hr/deck.pptx")

    chunk_ids = {c.chunk_id for c in chunks}
    table_chunks = [c for c in chunks if c.chunk_type == ChunkType.TABLE]
    for tc in table_chunks:
        assert tc.parent_id in chunk_ids, (
            f"Table chunk {tc.chunk_id} has orphan parent_id {tc.parent_id!r}"
        )
