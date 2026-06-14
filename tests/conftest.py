"""
Shared fixtures and dummy data for the ingestion pipeline test suite.

All Azure SDK clients are mocked — no real Azure resources are contacted.
Dummy values use realistic Azure resource name patterns and valid UUID formats.
"""
from __future__ import annotations

import io
from unittest.mock import MagicMock, patch
from uuid import uuid4

import pytest

# ── Dummy data constants ───────────────────────────────────────────────────────

DUMMY_TENANT_ID    = "aaaaaaaa-bbbb-cccc-dddd-eeeeeeeeeeee"
DUMMY_CLIENT_ID    = "11111111-2222-3333-4444-555555555555"
DUMMY_SITE_ID      = "ironman.sharepoint.com,aabbccdd-1122-3344-5566-778899aabbcc,ddee1122-3344-5566-7788-99aabbccddee"
DUMMY_DRIVE_ID     = "b!AbCdEfGhIjKlMnOpQrStUvWxYz1234567890ABCDEFGHIJ"
DUMMY_ITEM_ID      = "01ABCDEFGHIJKLMNOPQRSTUVWXYZ1234567"
DUMMY_SHA256       = "a" * 64  # valid 64-char hex
DUMMY_STORAGE_ACCT = "mystorageaccount"
DUMMY_SEARCH_EP    = "https://my-search.search.windows.net"
DUMMY_SB_NAMESPACE = "my-servicebus.servicebus.windows.net"
DUMMY_DOC_URL      = "https://ironman.sharepoint.com/sites/HR/Documents/policy.pdf"


def pytest_configure(config):
    config.addinivalue_line("markers", "asyncio: mark test as async")


# ── Environment fixture ────────────────────────────────────────────────────────

@pytest.fixture()
def env_vars():
    """Patch os.environ with all required env vars using DUMMY_* values."""
    env = {
        "AZURE_FOUNDRY_PROJECT_ENDPOINT": "https://my-foundry.api.azureml.ms",
        "AZURE_STORAGE_ACCOUNT_NAME":     DUMMY_STORAGE_ACCT,
        "AZURE_SEARCH_ENDPOINT":          DUMMY_SEARCH_EP,
        "AZURE_SEARCH_API_KEY":           "fake-search-api-key-32-chars-min",
        "AZURE_SERVICE_BUS_NAMESPACE":    DUMMY_SB_NAMESPACE,
        "SHAREPOINT_TENANT_ID":           DUMMY_TENANT_ID,
        "SHAREPOINT_CLIENT_ID":           DUMMY_CLIENT_ID,
        "SHAREPOINT_CLIENT_SECRET":       "super-secret-sharepoint-value-here",
        "SHAREPOINT_WEBHOOK_SECRET":      "somesecret-32chars-minimum-value",
    }
    with patch.dict("os.environ", env):
        yield env


# ── File byte fixtures ─────────────────────────────────────────────────────────

@pytest.fixture()
def sample_pdf_bytes() -> bytes:
    """Minimal valid PDF bytes (hand-crafted, no external dependency)."""
    # Minimal single-page PDF with one text object
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"
        b"4 0 obj\n<< /Length 44 >>\nstream\n"
        b"BT /F1 12 Tf 100 700 Td (Hello World) Tj ET\n"
        b"endstream\nendobj\n"
        b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
        b"xref\n0 6\n0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"0000000266 00000 n \n"
        b"0000000360 00000 n \n"
        b"trailer\n<< /Size 6 /Root 1 0 R >>\n"
        b"startxref\n441\n%%EOF\n"
    )
    return pdf


@pytest.fixture()
def sample_docx_bytes() -> bytes:
    """Minimal valid DOCX bytes created with python-docx in-memory."""
    import docx

    doc = docx.Document()
    doc.add_heading("Main Policy Heading", level=1)
    doc.add_paragraph(
        "This is a standard policy paragraph with enough content to trigger LLM cleaning."
    )
    tbl = doc.add_table(rows=2, cols=3)
    tbl.cell(0, 0).text = "Header A"
    tbl.cell(0, 1).text = "Header B"
    tbl.cell(0, 2).text = "Header C"
    tbl.cell(1, 0).text = "Value 1"
    tbl.cell(1, 1).text = "Value 2"
    tbl.cell(1, 2).text = "Value 3"

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture()
def sample_xlsx_bytes() -> bytes:
    """Minimal valid XLSX bytes created with openpyxl in-memory."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Department", "Salary"])
    ws.append(["Alice", "HR", 75000])
    ws.append(["Bob", "IT", 85000])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture()
def sample_pptx_bytes() -> bytes:
    """Minimal valid PPTX bytes created with python-pptx in-memory."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout

    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Q1 Results"
    slide1.placeholders[1].text = "Revenue increased by 15% in Q1 2024."

    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "HR Policy Update"
    slide2.placeholders[1].text = "New remote work policy effective from March 2024."

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Domain model fixture ───────────────────────────────────────────────────────

@pytest.fixture()
def raw_chunk_factory():
    """
    Returns a factory function that creates RawChunk instances with sensible
    defaults, allowing field-level overrides via keyword arguments.
    """
    from shared.models import ChunkType, RawChunk

    def make_chunk(**overrides) -> RawChunk:
        defaults = dict(
            chunk_id           = str(uuid4()),
            parent_id          = str(uuid4()),
            chunk_type         = ChunkType.PARAGRAPH,
            domain             = "hr",
            doc_name           = "policy.pdf",
            source             = "policy.pdf",
            doc_url            = DUMMY_DOC_URL,
            file_type          = "pdf",
            blob_path          = "hr/policy.pdf",
            ingested_at        = "2024-01-15T10:00:00+00:00",
            page_number        = 1,
            title              = "HR Policy 2024",
            section_heading    = "Leave Policy",
            section_subheading = "",
            content            = "Employees are entitled to 20 days of annual leave.",
            table_raw          = "",
            file_sha256        = DUMMY_SHA256,
            is_deleted         = False,
        )
        defaults.update(overrides)
        return RawChunk(**defaults)

    return make_chunk


# ── Azure client mocks ─────────────────────────────────────────────────────────

@pytest.fixture()
def mock_openai_client() -> MagicMock:
    """
    MagicMock mimicking AzureOpenAI.
    chat.completions.create returns 'Cleaned text output'.
    embeddings.create returns a 1536-dim vector.
    """
    client = MagicMock()

    # chat completions
    chat_choice = MagicMock()
    chat_choice.message.content = "Cleaned text output"
    chat_resp = MagicMock()
    chat_resp.choices = [chat_choice]
    client.chat.completions.create.return_value = chat_resp

    # embeddings
    emb_data = MagicMock()
    emb_data.embedding = [0.1] * 1536
    emb_resp = MagicMock()
    emb_resp.data = [emb_data]
    client.embeddings.create.return_value = emb_resp

    return client


@pytest.fixture()
def mock_search_client() -> MagicMock:
    """
    MagicMock mimicking azure.search.documents.SearchClient.
    search() returns empty iterator by default.
    upload_documents() returns one succeeded result.
    """
    client = MagicMock()

    client.search.return_value = iter([])

    succeeded_result = MagicMock()
    succeeded_result.succeeded = True
    client.upload_documents.return_value = [succeeded_result]
    client.delete_documents.return_value = [succeeded_result]

    return client
